import os
# from PyPDF2 import PdfReader
import fitz 
# PyMuPDF
from pptx import Presentation
from db_utils import save_slide
import re

def parse_document(file_path, document_id):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pdf':
        return parse_pdf(file_path, document_id)
    elif ext == '.pptx':
        return parse_pptx(file_path, document_id)
    else:
        raise ValueError("Unsupported file type")

def parse_pdf(file_path, document_id):
    try:
        # reader = PdfReader(file_path)
        doc = fitz.open(file_path)
    except Exception as e:
        raise ValueError(f"Error opening PDF file: {str(e)}")
        
    saved_slides = []
    # for i, page in enumerate(reader.pages):
    for i, page in enumerate(doc):
        try:
            # text = page.extract_text() or ''
            # raw_text = page.extract_text() or ''
            raw_text = page.get_text() or ''
            cleaned_text = clean_extracted_text(raw_text)

            slide_data = save_slide(
                slide_number=i + 1,
                title=f"PDF Page {i + 1}",
                # content=text.strip(),
                content=cleaned_text.strip(),
                document_id=document_id
            )
            saved_slides.append(slide_data)
        except Exception as e:
            print(f"⚠️ Error processing PDF page {i+1}: {e}")
            continue
    return saved_slides

def parse_pptx(file_path, document_id):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise ValueError(f"Error opening PPTX file: {str(e)}")
    
    saved_slides = []
    for i, slide in enumerate(prs.slides):
        try:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

            content = "\n".join(texts)
            slide_data = save_slide(
                slide_number=i + 1,
                title=f"PPT Slide {i + 1}",
                content=content.strip(),
                document_id=document_id
            )
            saved_slides.append(slide_data)
        except Exception as e:
            print(f"⚠️ Error processing PPT slide {i+1}: {e}")
            continue
    return saved_slides

def clean_extracted_text(text):
    """
    Attempt to improve spacing in extracted text.
    You can evolve this logic based on your dataset.
    """
    # Fix stuck lowercase-uppercase like "experienceIn"
    text = re.sub(r'(?<=[a-z])(?=[A-Z])', ' ', text)

    # Fix emails/URLs that get space-inserted (optional)
    text = re.sub(r'\s+([@.])\s+', r'\1', text)

    # Remove excess whitespace
    text = re.sub(r'\s{2,}', ' ', text)

    # Normalize newlines
    text = text.replace('\n', ' ').strip()

    text = re.sub(r'([a-z])([A-Z][a-z])', r'\1 \2', text)

    # Instead of replacing \n with space, collapse multiple into one line break
    text = re.sub(r'\n+', '\n', text)

    text = re.sub(r'\n{2,}', '\n', text)  # Keep single newlines

    return text
